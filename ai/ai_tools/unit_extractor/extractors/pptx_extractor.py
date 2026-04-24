from .base import ExtractorBase
from pptx import Presentation
from domain import TextData

class PptxExtractor(ExtractorBase):
    """
    Extracts pptx files
    """
    def _extract(self, path=None) -> list:
        
        prs = Presentation(path)
        
        slide_docs = []
        for i, slide in enumerate(prs.slides, start=1):
            
            text = self._get_text(slide)
            notes = self._get_notes(slide)

            data = TextData(path=path, text=f'CONTENT: {text}\n NOTES: {notes}\n')
            data.metadata['slide_num'] = f'{i}'
            
            slide_docs.append(data)
            
        return slide_docs
    
    def _get_text(self, slide) -> str:
        """Retrieves all text from shapes in order"""
        collected_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                collected_text.append(shape.text)
        return "".join(collected_text)
 
    def _get_notes(self, slide) -> str:
        """Retrieves the speaker notes/script"""
        if slide.has_notes_slide:
            frame = slide.notes_slide.notes_text_frame
            if frame is not None:
                return frame.text
        return ""

if __name__ == "__main__":
    
    test_path = r"C:\All University Materials\Project\ICT304-project\ai\ai_tools\data_extractor\data\raw\ICT283_contents\Topic04\Content\Lec-10.pptx"
    
    for i in PptxExtractor().extract(test_path):
        print(i.metadata)
                
